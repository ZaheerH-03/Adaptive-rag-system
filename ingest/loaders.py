from pypdf import PdfReader
from .schema import DocUnit
from pathlib import Path
from typing import List
from pptx import Presentation
from docx import Document as DocxDocument
def load_pdf_units(path: Path) -> List[DocUnit]:
    reader = PdfReader(str(path))
    units: List[DocUnit] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        text = text.strip()
        if not text:
            continue

        units.append(
            DocUnit(
                text=text,
                filename=path.name,
                file_type = 'pdf',
                page_num= i+1
            )
        )
    return units

def load_pptx_units(path: Path) -> List[DocUnit]:
    prs = Presentation(str(path))
    units: List[DocUnit] = []
    for i,slide in enumerate(prs.slides):
        parts = []
        title_text = None

        if slide.shapes.title is not None and hasattr(slide.shapes.title,"text"):
            title_text = slide.shapes.title.text.strip()
            if title_text:
                parts.append(title_text)
        for shape in slide.shapes:
            if shape is slide.shapes.title:
                continue
            if hasattr(shape,"text"):
                t = shape.text.strip()
                if t:
                    parts.append(t)
        full_text = "\n".join(parts).strip()
        if not full_text:
            continue
        units.append(
            DocUnit(
                text=full_text,
                filename=path.name,
                file_type='pptx',
                slide_num= i+1,
                section_title=title_text,
            )
        )
    return units

def load_docx_units(path: Path) -> List[DocUnit]:
    doc = DocxDocument(str(path))
    units: List[DocUnit] = []

    current_title = None
    buffer = []

    def flush():
        nonlocal buffer, current_title
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        if not text:
            buffer = []
            return
        units.append(
            DocUnit(
                text=text,
                filename=path.name,
                file_type="docx",
                section_title=current_title,
            )
        )
        buffer = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            flush()
            current_title = text
        else:
            buffer.append(text)

    flush()
    return units


def load_units_for_file(path: Path) -> List[DocUnit]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return load_pdf_units(path)
    elif ext == ".pptx":
        return load_pptx_units(path)
    elif ext == ".docx":
        return load_docx_units(path)
    else:
        return []